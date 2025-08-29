import os
import json
import httpx
import pandas as pd
import requests
from bs4 import BeautifulSoup
from typing import List, Dict
from openai import OpenAI
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from datetime import datetime

# Credentials

api_key = "xxx"
search_engine_id = "yyy"
openai_api_key = "ZZZ"

# Extracting Information from the Web

query = input("Enter the topic on which the Slides has to be prepared:")

def googleSearch(api_key, search_engine_id, query, **params):
    base_url = "https://www.googleapis.com/customsearch/v1"
    params = {
        'key': api_key,
        'cx': search_engine_id,
        'q': query,
        **params
    }
    response = httpx.get(base_url, params=params)
    response.raise_for_status()
    return response.json()

results = []
response = googleSearch(api_key, search_engine_id, query)
results.extend(response.get('items', []))

df = pd.json_normalize(results)
data = df[["title", "link"]]
data.head()

def extractText(url, max_chars=4000):
    try:
        resp = requests.get(url, timeout=10)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")

        # Remove scripts/styles
        for script in soup(["script", "style", "noscript"]):
            script.extract()

        text = " ".join(soup.stripped_strings)
        return text[:max_chars]  # clip to avoid too long input for LLM
    except Exception as e:
        print(f"Error fetching {url}: {e}")
        return ""

Content = [] # Stores the title, url and extracted text from the WebSources

for i in range(10):
    title = data["title"][i]
    url = data["link"][i]
    text = extractText(url)
    if text!="" and len(Content)<5:
        Content.append({"title": title,
                        "url": url,
                        "text": text})

# LLM Inclusion
def GenerateSLide_Outline(
    WebContent: List[Dict[str, str]],
    query: str):
    """
    WebContent: [{"title": str, "url": str, "text": str}]
    query: original topic/query string
    Returns a dict with 7-slide outline.
    """

    # Normalize & trim to control tokens
    compact_sources = []
    for c in WebContent[:5]:
        compact_sources.append({
            "title": c.get("title", "")[:200],
            "url": c.get("url", "")[:1000],
            "extract": (c.get("text", "") or "")[:2500]
        })


    client = OpenAI(api_key=openai_api_key)
    model = "gpt-4o-mini" # can be changed to the desired OpenAI model
    
    # Strict JSON schema for slides
    schema = {
        "name": "SlideDeck",
        "schema": {
            "type": "object",
            "properties": {
                "title": {"type": "string"},
                "overview": {"type": "array", "items": {"type": "string"}, "minItems": 4, "maxItems": 6},
                "slides": {
                    "type": "array",
                    "minItems": 4, "maxItems": 4,
                    "items": {
                        "type": "object",
                        "properties": {
                            "heading": {"type": "string"},
                            "bullets": {"type": "array", "items": {"type": "string"}, "minItems": 3, "maxItems": 6}
                        },
                        "required": ["heading", "bullets"],
                        "additionalProperties": False
                    }
                },
                "takeaways": {"type": "array", "items": {"type": "string"}, "minItems": 3, "maxItems": 5},
                "sources": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {"title": {"type": "string"}, "url": {"type": "string"}},
                        "required": ["title", "url"],
                        "additionalProperties": False
                    },
                    "minItems": 5, "maxItems": 10
                }
            },
            "required": ["title", "overview", "slides", "takeaways", "sources"],
            "additionalProperties": False
        },
        "strict": True
    }

    

    system_msg = (
        "You are SlideSmith, an expert research assistant that synthesizes a crisp, factual presentation outline. "
        "Write in concise bullets, avoid marketing fluff, and prefer concrete facts with dates. "
        "All content must be self-contained, objective, and neutral. No first-person voice. "
        "Combine your world knowledge with the provided sources; when using claims from sources, prefer those facts and "
        "reflect them accurately. Do not invent URLs."
    )

    user_text = (
        f"Topic: {query}\n\n"
        "Produce a 7-slide outline:\n"
        "  Slide 1: Title (just the title string)\n"
        "  Slide 2: Overview (4–6 bullets)\n"
        "  Slides 3–6: Four key sections (each has a heading and 3–6 bullets)\n"
        "  Slide 7: Takeaways (3–5 bullets)\n\n"
        "Rules:\n"
        " - Use dates and numbers when available.\n"
        " - Keep each bullet under 25 words.\n"
        " - Do not invent URLs; cite from provided results where possible.\n\n"
        f"WebSources: {compact_sources}"
    )

    response = client.chat.completions.create(
        model=model,
        messages=[
            {"role": "system", "content": system_msg},
            {"role": "user", "content": user_text}
        ],
        response_format={
        "type": "json_schema",  # required
        "json_schema": schema
        },
        temperature=0.3,
        max_tokens=1200
    )
    
    # Parse the response
    return json.loads(response.choices[0].message.content)

slidedeck = GenerateSLide_Outline(Content, query)

# Creating PowerPoint Presentation

def create_ppt_from_slidedeck(slidedeck, filename):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]  # title slide
    bullet_slide_layout = prs.slide_layouts[1] # title + content

    # --- Helper function to style heading ---
    def style_heading(shape):
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.size = Pt(36)
        run.font.bold = True
        run.font.color.rgb = RGBColor(0x2E, 0x74, 0xB5)  # nice blue

    # --- Helper function to style body text ---
    def style_body(tf):
        for p in tf.paragraphs:
            for run in p.runs:
                run.font.size = Pt(20)
        tf.word_wrap = True

    # --- Helper: add footer with date + page number ---
    def add_footer(slide, page_num, total_pages):
        left = Inches(0.3)
        top = Inches(6.8)
        width = Inches(4)
        height = Inches(0.3)

        # Date (left side)
        textbox = slide.shapes.add_textbox(left, top, width, height)
        textbox.text = datetime.today().strftime("%d %b %Y")
        textbox.text_frame.paragraphs[0].font.size = Pt(12)

        # Page number (right side)
        textbox = slide.shapes.add_textbox(Inches(9), top, width, height)
        textbox.text = f"{page_num}/{total_pages}"
        textbox.text_frame.paragraphs[0].font.size = Pt(12)

    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = slidedeck["title"]
    slide.placeholders[1].text = ""  # optional subtitle
    style_heading(slide.shapes.title)

    # Collect all slides except title for footer count
    sections = [
        ("Overview", slidedeck["overview"]),
        *[(s["heading"], s["bullets"]) for s in slidedeck["slides"]],
        ("Key Takeaways", slidedeck["takeaways"]),
        ("Sources", [f"{s['title']} - {s['url']}" for s in slidedeck["sources"]])
    ]

    total_pages = len(sections)

    # Slides
    for i, (heading, bullets) in enumerate(sections, start=1):
        slide = prs.slides.add_slide(bullet_slide_layout)
        slide.shapes.title.text = heading
        style_heading(slide.shapes.title)

        tf = slide.placeholders[1].text_frame
        tf.clear()
        for bullet in bullets:
            p = tf.add_paragraph()
            p.text = bullet
        style_body(tf)

        # Footer
        add_footer(slide, i, total_pages)

    prs.save(filename)
    return filename

create_ppt_from_slidedeck(slidedeck, f"{query}.pptx")

